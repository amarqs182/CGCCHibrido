from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Cores
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MED = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
GREEN = RGBColor(0x00, 0xe6, 0x76)
ORANGE = RGBColor(0xff, 0x8c, 0x00)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

# ==================== SLIDE 1: TITULO ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 2.8, 11.3)

add_textbox(slide, 1, 1.5, 11.3, 1.5, "Diorama 3D", 48, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.0, 11.3, 1, "Visualizador 3D Completo - Computacao Grafica 2026/1", 24, WHITE, False, PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.2, 11.3, 0.5, "Unisinos - Ciencia da Computacao (Hibrido)", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.0, 11.3, 0.5, "Aluno: amarqs182", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ==================== SLIDE 2: OBJETIVO ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Objetivo", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

add_bullet_list(slide, 0.8, 1.8, 11.7, 5, [
    "Integrar todos os topicos do semestre em uma unica aplicacao",
    "Demonstrar dominio do pipeline grafico completo",
    "Visualizador funcional com interacao em tempo real",
    "",
    "Topicos integrados:",
    "   OpenGL 4.x  |  Camera FPS  |  Iluminacao Phong",
    "   Curvas de Bezier  |  Texturas  |  Carregamento .obj",
    "   Interface ImGui  |  Materiais .mtl"
], 20, WHITE)

# ==================== SLIDE 3: ARQUITETURA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Arquitetura", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

add_bullet_list(slide, 0.8, 1.8, 5.5, 5, [
    "MainDiorama.cpp - Aplicacao principal",
    "Camera.cpp - Camera FPS (WASD + mouse)",
    "ObjParser.cpp - Parser .obj/.mtl",
    "Trajectory.cpp - Curvas de Bezier",
    "SceneParser.cpp - Leitura de cena JSON"
], 18, WHITE)

add_bullet_list(slide, 7, 1.8, 5.5, 5, [
    "Dependencias (FetchContent):",
    "   GLFW 3.4 - Janela e input",
    "   GLM - Matrizes e math",
    "   stb_image - Carregamento de imagens",
    "   ImGui 1.91.8 - Interface grafica",
    "   GLAD - OpenGL loader"
], 18, LIGHT_GRAY)

# ==================== SLIDE 4: CENA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Cena do Diorama", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

add_bullet_list(slide, 0.8, 1.8, 5.5, 5, [
    "5 Objetos:",
    "   Cubo Principal (vermelho)",
    "   Esfera Texturizada (pixelWall)",
    "   Piramide Azul",
    "   Chao (cinza)",
    "   Suzanne - modelo .obj (macaco)"
], 18, WHITE)

add_bullet_list(slide, 7, 1.8, 5.5, 5, [
    "3 Fontes de Luz (Phong):",
    "   Key Light - luz principal (branca)",
    "   Fill Light - luz secundaria (azulada)",
    "   Back Light - luz traseira (quente)",
    "",
    "Arquivo de configuracao: scene.json"
], 18, LIGHT_GRAY)

# ==================== SLIDE 5: FUNCIONALIDADES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Funcionalidades", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

# Coluna 1
add_textbox(slide, 0.8, 1.8, 5.5, 0.5, "Camera FPS", 22, GREEN, True)
add_bullet_list(slide, 0.8, 2.3, 5.5, 1.5, [
    "WASD - Mover",
    "Q/E - Subir/Descer",
    "Mouse - Olhar",
    "Scroll - Zoom"
], 16, WHITE)

add_textbox(slide, 0.8, 4.0, 5.5, 0.5, "Transformacoes", 22, GREEN, True)
add_bullet_list(slide, 0.8, 4.5, 5.5, 1.5, [
    "TAB - Trocar objeto",
    "IJKL - Transladar",
    "U/O - Rotacionar",
    "+/- - Escalar"
], 16, WHITE)

# Coluna 2
add_textbox(slide, 7, 1.8, 5.5, 0.5, "Iluminacao Phong", 22, GREEN, True)
add_bullet_list(slide, 7, 2.3, 5.5, 1.5, [
    "F1/F2/F3 - Lig/deslig luz 1/2/3",
    "F4 - Lig/deslig todas",
    "ImGui - Editar posicao/cor/intensidade"
], 16, WHITE)

add_textbox(slide, 7, 4.0, 5.5, 0.5, "Texturas e Materiais", 22, GREEN, True)
add_bullet_list(slide, 7, 4.5, 5.5, 1.5, [
    "Carregamento automatico .mtl",
    "Troca manual via ImGui",
    "Edicao de Ka, Kd, Ks, shininess"
], 16, WHITE)

add_textbox(slide, 0.8, 6.0, 11.7, 0.5, "Animacao: SPACE (ativar) | B (Bezier) | 1-5 (pontos) | R (resetar) | UP/DOWN (velocidade)", 16, ORANGE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 6: SHADER PHONG ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Iluminacao Phong - Fragment Shader", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 5)

add_textbox(slide, 0.8, 1.8, 11.7, 0.5, "Modelo de iluminacao com 3 componentes:", 20, WHITE, True)

add_bullet_list(slide, 0.8, 2.5, 11.7, 4.5, [
    "Ambient:  Ka * lightColor * lightIntensity",
    "",
    "Diffuse:  Kd * max(dot(N, L), 0.0) * lightColor * lightIntensity",
    "",
    "Specular: Ks * pow(max(dot(V, R), 0.0), shininess) * lightColor * lightIntensity",
    "",
    "Atenuacao: 1.0 / (Kc + Kl*d + Kq*d^2)",
    "",
    "Resultado = (ambient + diffuse) * objectColor + specular"
], 18, LIGHT_GRAY)

# ==================== SLIDE 7: DESAFIOS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Desafios Tecnicos", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

add_bullet_list(slide, 0.8, 1.8, 11.7, 5, [
    "Dangling pointer: SceneObject local era destruido antes do render loop",
    "   Solucao: armazenar por valor (SceneObject config) em vez de ponteiro",
    "",
    "Shaders #version 450 vs OpenGL 4.1 context",
    "   Solucao: alterar para #version 410 para compatibilidade",
    "",
    "Shallow clone impedia push para GitHub",
    "   Solucao: git fetch --unshallow antes do push",
    "",
    "Parser .obj com suporte a multi-materiais e texturas"
], 18, WHITE)

# ==================== SLIDE 8: REFERENCIAS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, 0.8, 0.5, 11.7, 0.8, "Referencias", 36, ACCENT, True)
add_accent_line(slide, 0.8, 1.3, 3)

add_bullet_list(slide, 0.8, 1.8, 11.7, 5, [
    "LearnOpenGL - https://learnopengl.com/",
    "OpenGL Documentation - https://docs.gl/",
    "GLM Documentation - https://glm.g-truc.net/",
    "ImGui Documentation - https://github.com/ocornut/imgui",
    "GLAD Generator - https://glad.dav1d.de/",
    "CMake FetchContent - https://cmake.org/cmake/help/latest/module/FetchContent.html",
    "Suzanne model - Blender (开源)",
    "Texturas - https://ambientcg.com/"
], 18, LIGHT_GRAY)

# ==================== SLIDE 9: OBRIGADO ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, 1, 3.2, 11.3)

add_textbox(slide, 1, 2.0, 11.3, 1.5, "Obrigado!", 60, ACCENT, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.5, 11.3, 1, "github.com/amarqs182/CGCCHibrido", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.5, 11.3, 1, "Perguntas?", 28, WHITE, False, PP_ALIGN.CENTER)

# Salvar
output_path = "/home/asm/unisinos/computação grafica/Vault/04-Courses/Computação Gráfica/CGCCHibrido/apresentacao.pptx"
prs.save(output_path)
print(f"Apresentacao salva em: {output_path}")
